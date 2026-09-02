#!/usr/bin/env python3
"""Generate RSAIRF.pptx from the Really Simple AI Risk Framework content.

Mirrors the scroll-deck in presentation.html, slide for slide (s1-s14).
Run: python3 build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Brand palette (from presentation.html :root) -------------------------
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
PAPER_RAISED = RGBColor(0xF8, 0xF8, 0xF8)
INK = RGBColor(0x1B, 0x1F, 0x23)
INK_SOFT = RGBColor(0x45, 0x45, 0x45)
INK_FAINT = RGBColor(0x99, 0x99, 0x99)
LINE = RGBColor(0xE0, 0xE0, 0xE0)
ACCENT = RGBColor(0x3E, 0x89, 0x14)          # green — primary accent
ACCENT_INK = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_SECONDARY = RGBColor(0x05, 0x66, 0x8D)  # teal — links
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# POV matrix scale (likelihood x impact heat-map)
S_LOW = RGBColor(0x38, 0x7A, 0x13)
S_MED = RGBColor(0xFF, 0xF1, 0x2C)
S_MED_INK = RGBColor(0x45, 0x45, 0x45)
S_HIGH = RGBColor(0xD6, 0x6F, 0x2C)
S_CRIT = RGBColor(0xCC, 0x23, 0x2A)

# Dark title/closing surface (body bg approximated as near-black, matching --ink)
DARK_BG = RGBColor(0x1B, 0x1F, 0x23)

FONT_HEAD = "Exo 2"        # headings / kickers
FONT_SUB = "Roboto Slab"   # ledes / subs
FONT_BODY = "Roboto"       # body text
FONT_MONO = "Roboto Mono"  # numbers / tags / mono bits

# 16:9 canvas
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def add(a, b):
    return Emu(int(a) + int(b))


def rect(s, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE, dashed=False):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
        if dashed:
            ln = sp.line._get_or_add_ln()
            ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return sp


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb, tf


def setpar(p, text, size, color=INK, bold=False, italic=False, align=PP_ALIGN.LEFT,
           font=FONT_BODY, space_after=None, space_before=None, line=None):
    p.text = text
    p.alignment = align
    if space_after is not None:
        p.space_after = space_after
    if space_before is not None:
        p.space_before = space_before
    if line is not None:
        p.line_spacing = line
    r = p.runs[0]
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return p


def add_runs(p, parts, size, align=PP_ALIGN.LEFT, line=None, space_after=None, font=FONT_BODY):
    """parts: list of (text, color, bold)."""
    p.alignment = align
    if line is not None:
        p.line_spacing = line
    if space_after is not None:
        p.space_after = space_after
    for text, color, bold in parts:
        r = p.add_run()
        r.text = text
        r.font.size = size
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
    return p


def bg(s, color):
    r = rect(s, 0, 0, SW, SH, fill=color)
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return r


def kicker(s, text, x=Inches(0.7), y=Inches(0.55)):
    tb, tf = textbox(s, x, y, Inches(6), Inches(0.35))
    setpar(tf.paragraphs[0], text.upper(), Pt(12), color=ACCENT, bold=True, font=FONT_HEAD)
    return tb


def title(s, text, y=Inches(0.95), size=Pt(30), w=Inches(11.9)):
    tb, tf = textbox(s, Inches(0.7), y, w, Inches(0.9))
    setpar(tf.paragraphs[0], text, size, color=INK, bold=False, font=FONT_HEAD, line=1.05)
    return tb


def footnote(s, text, y=Inches(6.95)):
    tb, tf = textbox(s, Inches(0.7), y, Inches(11.9), Inches(0.45))
    setpar(tf.paragraphs[0], text, Pt(10.5), color=INK_FAINT, line=1.1)


def counter(s, n, total=14):
    tb, tf = textbox(s, Inches(0.5), Inches(6.95), Inches(1.5), Inches(0.35))
    add_runs(tf.paragraphs[0], [(f"{n:02d}", INK, True), (f" / {total}", INK_SOFT, False)],
             Pt(10.5), font=FONT_MONO)


def wordmark(s):
    tb, tf = textbox(s, Inches(11.6), Inches(6.95), Inches(1.2), Inches(0.35))
    setpar(tf.paragraphs[0], "RSAIRF", Pt(10), color=INK_FAINT, align=PP_ALIGN.RIGHT, font=FONT_MONO)


def base_slide(n, wide=False):
    s = slide()
    bg(s, PAPER)
    counter(s, n)
    wordmark(s)
    return s


def stat_tile(s, x, y, w, h, num, label):
    rect(s, x, y, w, h, fill=None, line=LINE, line_w=Pt(0.75))
    tb, tf = textbox(s, x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), h - Inches(0.2))
    setpar(tf.paragraphs[0], num, Pt(26), color=INK, bold=True, font=FONT_MONO, space_after=Pt(3))
    p2 = tf.add_paragraph()
    setpar(p2, label, Pt(11), color=INK_SOFT, line=1.1)


def table_grid(s, x, y, col_ws, rows, row_h=Inches(0.5), header=True, font_size=Pt(11.5)):
    ry = y
    for ri, row in enumerate(rows):
        cx = x
        is_header = header and ri == 0
        for ci, val in enumerate(row):
            w = col_ws[ci]
            fillc = PAPER
            txtc = INK_SOFT if not is_header else INK_SOFT
            bold = False
            if is_header:
                fillc = None
            elif ri % 2 == 0:
                fillc = PAPER_RAISED
            rect(s, cx, ry, w, row_h, fill=fillc, line=LINE if not is_header else None,
                 line_w=Pt(0.5))
            if is_header:
                # bottom border only, like the HTML table's th
                bline = rect(s, cx, add(ry, row_h) - Emu(int(Pt(1))), w, Emu(int(Pt(1))), fill=INK_SOFT)
            tb, tf = textbox(s, cx + Inches(0.1), ry, w - Inches(0.15), row_h,
                             anchor=MSO_ANCHOR.MIDDLE)
            if is_header:
                setpar(tf.paragraphs[0], val, Pt(9.5), color=INK_SOFT, bold=False,
                       font=FONT_MONO, line=1.0)
            else:
                col0 = (ci == 0)
                setpar(tf.paragraphs[0], val, font_size, color=INK if not col0 else INK,
                       bold=False, font=FONT_MONO if col0 else FONT_BODY, line=1.05)
            cx = add(cx, w)
        ry = add(ry, row_h)
    return ry


# =========================================================================
# S1. TITLE
# =========================================================================
s = slide()
bg(s, DARK_BG)
rect(s, Inches(0.7), Inches(0.75), Inches(0.14), Inches(0.14), fill=ACCENT)
tb, tf = textbox(s, Inches(0.95), Inches(0.68), Inches(5), Inches(0.35), anchor=MSO_ANCHOR.MIDDLE)
setpar(tf.paragraphs[0], "AI RISK GOVERNANCE", Pt(11), color=RGBColor(0xC7, 0xC7, 0xC7),
       bold=True, font=FONT_HEAD)

tb, tf = textbox(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(2.6))
setpar(tf.paragraphs[0], "Really Simple", Pt(58), color=WHITE, bold=False, font=FONT_HEAD,
       space_after=Pt(0), line=1.0)
p = tf.add_paragraph()
setpar(p, "AI Risk Framework", Pt(58), color=WHITE, bold=False, font=FONT_HEAD, line=1.0)

tb, tf = textbox(s, Inches(0.7), Inches(4.15), Inches(9.5), Inches(1.1))
setpar(tf.paragraphs[0],
       "A shared register for naming, scoring, and tracking what can go wrong when you build, "
       "ship, or run AI — thirty risks, four lifecycle stages, one ID space.",
       Pt(17), color=RGBColor(0xDD, 0xDD, 0xDD), font=FONT_SUB, line=1.4)

tb, tf = textbox(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5))
add_runs(tf.paragraphs[0], [
    ("RSAIRF · Open Ethics Initiative     ", RGBColor(0xAA, 0xAA, 0xAA), False),
    ("github.com/OpenEthicsAI/RSAIRF     ", RGBColor(0xAA, 0xAA, 0xAA), False),
    ("Licensed CC BY 4.0", RGBColor(0xAA, 0xAA, 0xAA), False),
], Pt(10.5))

tb, tf = textbox(s, Inches(8.3), Inches(0.2), Inches(4.8), Inches(2.0), anchor=MSO_ANCHOR.MIDDLE)
setpar(tf.paragraphs[0], "01–30", Pt(90), color=RGBColor(0x2A, 0x2E, 0x33), bold=True,
       font=FONT_MONO, align=PP_ALIGN.RIGHT)

# =========================================================================
# S2. WHY IT EXISTS
# =========================================================================
s = base_slide(2)
kicker(s, "Why it exists")
tb, tf = textbox(s, Inches(0.7), Inches(1.05), Inches(9.5), Inches(1.3))
setpar(tf.paragraphs[0], "Every organization deploys and uses AI.", Pt(26), color=INK,
       font=FONT_HEAD, line=1.15, space_after=Pt(0))
p = tf.add_paragraph()
setpar(p, "Few build it. Almost none monitor it well.", Pt(26), color=INK, font=FONT_HEAD,
       line=1.15)

tb, tf = textbox(s, Inches(0.7), Inches(2.4), Inches(11.4), Inches(1.9))
setpar(tf.paragraphs[0],
       "The AI Risk Register is curated from public sources and the internal incident learnings "
       "of real client engagements — not written from theory. It exists because the biggest "
       "exposures usually aren’t in the model itself. They’re in the gaps around it: tools "
       "adopted without approval, and monitoring that was never built.",
       Pt(15), color=INK_SOFT, font=FONT_SUB, line=1.5)

stats = [("30", "registered risks"), ("4", "lifecycle stages"), ("AIR01–30", "traceable IDs")]
sx = Inches(0.7)
sw_ = Inches(3.9)
sgap = Inches(0.02)
sy = Inches(4.6)
sh_ = Inches(1.3)
for num, lbl in stats:
    stat_tile(s, sx, sy, sw_, sh_, num, lbl)
    sx = add(add(sx, sw_), sgap)

# =========================================================================
# S3. THE LIFECYCLE
# =========================================================================
s = base_slide(3)
kicker(s, "The lifecycle")
title(s, "Every initiative moves through four stages", y=Inches(1.05), size=Pt(28))

stage_names = [("Development", "Model Development"),
               ("Deployment", "Model / Component Deployment"),
               ("Use", "AI Use")]
bx = Inches(0.7)
bw = Inches(3.55)
by = Inches(2.15)
bh = Inches(1.05)
bgap = Inches(0.55)
for i, (name, role) in enumerate(stage_names):
    rect(s, bx, by, bw, bh, fill=PAPER_RAISED, line=LINE, line_w=Pt(0.75))
    tb, tf = textbox(s, bx + Inches(0.2), by + Inches(0.12), bw - Inches(0.4), bh - Inches(0.2))
    setpar(tf.paragraphs[0], name, Pt(16), color=INK, bold=False, font=FONT_HEAD, space_after=Pt(3))
    p2 = tf.add_paragraph()
    setpar(p2, role, Pt(11), color=INK_SOFT)
    if i < 2:
        tb2, tf2 = textbox(s, add(bx, bw), by, bgap, bh, anchor=MSO_ANCHOR.MIDDLE)
        setpar(tf2.paragraphs[0], "→", Pt(20), color=INK_FAINT, align=PP_ALIGN.CENTER,
               font=FONT_MONO)
    bx = add(add(bx, bw), bgap)

# monitor band
my = add(by, add(bh, Inches(0.35)))
rect(s, Inches(0.7), my, Inches(11.9), Inches(0.65), fill=None, line=LINE, line_w=Pt(0.75), dashed=True)
tb, tf = textbox(s, Inches(0.7), my, Inches(11.9), Inches(0.65), anchor=MSO_ANCHOR.MIDDLE)
add_runs(tf.paragraphs[0], [
    ("Continuous Monitoring & Evaluation", INK, True),
    ("  —  runs across all three stages, and feeds findings back into any of them", INK_SOFT, False),
], Pt(12.5), align=PP_ALIGN.CENTER)

defs = [
    ("Development", "You design or materially change AI behavior — data selection, training, tuning, evaluation."),
    ("Deployment", "You don’t build the model, but you integrate, configure, and release AI components into production."),
    ("Use", "You consume a deployed AI solution and make decisions or take actions based on its outputs."),
    ("Monitoring", "You track risk, quality, and incidents over time — continuous, and running across every other stage."),
]
dx = Inches(0.7)
dw = Inches(2.87)
dgap = Inches(0.13)
dy = add(my, Inches(1.0))
for name, body in defs:
    tb, tf = textbox(s, dx, dy, dw, Inches(1.6))
    setpar(tf.paragraphs[0], name, Pt(13), color=ACCENT, bold=False, font=FONT_HEAD, space_after=Pt(4))
    p2 = tf.add_paragraph()
    setpar(p2, body, Pt(10.5), color=INK_SOFT, line=1.35)
    dx = add(add(dx, dw), dgap)

# =========================================================================
# S4. THE PROCESS
# =========================================================================
s = base_slide(4)
kicker(s, "The process")
title(s, "Four steps, repeated at every transition", y=Inches(1.05), size=Pt(28))

steps = [
    "Check which lifecycle stages apply to your initiative — run a discovery session with your "
    "technical team if it’s unclear.",
    "Map risks using register.csv. Stage tags show where each risk is typically introduced or "
    "first exploitable.",
    "For every applicable AIR## risk, record likelihood, impact, an owner, and a mitigation — and "
    "reference that ID in your project’s risk log or audit.",
    "Found a risk that isn’t represented? Open a repository issue — your contribution may help "
    "other teams too.",
]
sy = Inches(2.15)
for i, text in enumerate(steps, start=1):
    tb, tf = textbox(s, Inches(0.7), sy, Inches(0.6), Inches(0.7))
    setpar(tf.paragraphs[0], f"0{i}", Pt(13), color=ACCENT, font=FONT_MONO)
    tb2, tf2 = textbox(s, Inches(1.35), sy, Inches(10.7), Inches(0.7))
    setpar(tf2.paragraphs[0], text, Pt(13.5), color=INK, line=1.3)
    sy = add(sy, Inches(0.72))

rect(s, Inches(0.7), add(sy, Inches(0.1)), Inches(0.04), Inches(0.55), fill=ACCENT)
tb, tf = textbox(s, Inches(0.9), add(sy, Inches(0.08)), Inches(10.8), Inches(0.6))
setpar(tf.paragraphs[0],
       "Re-check the register at each stage transition — a risk that was out of scope earlier "
       "may now apply.", Pt(13.5), color=INK, italic=False, font=FONT_SUB, line=1.35)

tb, tf = textbox(s, Inches(0.7), add(sy, Inches(0.85)), Inches(10.8), Inches(0.4))
setpar(tf.paragraphs[0], "↻  treat this as a living list, not a one-time checklist", Pt(11.5),
       color=INK_SOFT)

# =========================================================================
# S5. SCORING · 1 OF 2 — LIKELIHOOD
# =========================================================================
s = base_slide(5)
kicker(s, "Scoring · 1 of 2")
title(s, "How probable it is", y=Inches(1.05), size=Pt(28))
tb, tf = textbox(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.7))
add_runs(tf.paragraphs[0], [
    ("For every applicable risk, rate ", INK_SOFT, False),
    ("likelihood", INK, True),
    (" — how probable it is that the risk materializes in your context.", INK_SOFT, False),
], Pt(14.5), font=FONT_SUB, line=1.4)

like_rows = [
    ("Score", "Likelihood", "Meaning"),
    ("1", "Rare", "Not expected during the initiative’s lifetime."),
    ("2", "Unlikely", "Could occur, but not expected."),
    ("3", "Possible", "May occur occasionally."),
    ("4", "Likely", "Expected to occur at some point."),
    ("5", "Almost certain", "Expected to occur often, or already observed."),
]
table_grid(s, Inches(0.7), Inches(2.75), [Inches(1.2), Inches(2.4), Inches(8.3)], like_rows,
           row_h=Inches(0.58))

# =========================================================================
# S6. SCORING · 2 OF 2 — IMPACT
# =========================================================================
s = base_slide(6)
kicker(s, "Scoring · 2 of 2")
title(s, "How severe, if it happens", y=Inches(1.05), size=Pt(28))
tb, tf = textbox(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.7))
add_runs(tf.paragraphs[0], [
    ("Then rate ", INK_SOFT, False),
    ("impact", INK, True),
    (" — the severity of the consequence if the risk does materialize.", INK_SOFT, False),
], Pt(14.5), font=FONT_SUB, line=1.4)

impact_rows = [
    ("Score", "Impact", "Meaning"),
    ("1", "Negligible", "Minimal effect; no material harm."),
    ("2", "Minor", "Limited harm, easily remediated."),
    ("3", "Moderate", "Noticeable harm, cost, or user impact."),
    ("4", "Major", "Significant financial, regulatory, or reputational harm."),
    ("5", "Severe", "Critical or potentially irreversible harm to people, rights, or the organization."),
]
end_y = table_grid(s, Inches(0.7), Inches(2.6), [Inches(1.2), Inches(2.4), Inches(8.3)],
                    impact_rows, row_h=Inches(0.55))

fb = rect(s, Inches(0.7), add(end_y, Inches(0.25)), Inches(6.2), Inches(0.65), fill=PAPER_RAISED,
          line=LINE, line_w=Pt(0.75))
tb, tf = textbox(s, Inches(0.9), add(end_y, Inches(0.25)), Inches(5.8), Inches(0.65),
                 anchor=MSO_ANCHOR.MIDDLE)
add_runs(tf.paragraphs[0], [
    ("Risk score = ", INK, False),
    ("Likelihood", ACCENT, True),
    (" × ", INK, False),
    ("Impact", ACCENT, True),
    ("   ·   range 1–25", INK, False),
], Pt(14), font=FONT_MONO)

# =========================================================================
# S7. THE MATRIX
# =========================================================================
s = base_slide(7)
kicker(s, "The matrix")
title(s, "Likelihood × impact, at a glance", y=Inches(1.05), size=Pt(28))
tb, tf = textbox(s, Inches(0.7), Inches(1.85), Inches(11), Inches(0.6))
setpar(tf.paragraphs[0],
       "The same 1–25 score, banded into four priority levels — consistent across "
       "initiatives, without heavy tooling.", Pt(14.5), color=INK_SOFT, font=FONT_SUB, line=1.4)

likelihoods = ["Rare", "Unlikely", "Possible", "Likely", "Almost certain"]


def band_color(v):
    if v <= 4:
        return S_LOW
    if v <= 9:
        return S_MED
    if v <= 14:
        return S_HIGH
    return S_CRIT


gx0 = Inches(2.1)
gy0 = Inches(3.1)
gcell = Inches(0.85)
rowhead_w = Inches(2.1)

tb, tf = textbox(s, gx0, Inches(2.75), Inches(4.3), Inches(0.3))
setpar(tf.paragraphs[0], "IMPACT →", Pt(9.5), color=INK_FAINT, font=FONT_MONO)

for c in range(5):
    cx = add(gx0, Emu(c * int(gcell)))
    tb, tf = textbox(s, cx, add(gy0, Inches(-0.32)), gcell, Inches(0.3), anchor=MSO_ANCHOR.BOTTOM)
    setpar(tf.paragraphs[0], str(c + 1), Pt(10.5), color=INK_SOFT, align=PP_ALIGN.CENTER,
           font=FONT_MONO)

for l in range(1, 6):
    ry = add(gy0, Emu((l - 1) * int(gcell)))
    tb, tf = textbox(s, Inches(0.7), ry, rowhead_w, gcell, anchor=MSO_ANCHOR.MIDDLE)
    setpar(tf.paragraphs[0], f"{l} · {likelihoods[l-1]}", Pt(10.5), color=INK_SOFT,
           align=PP_ALIGN.RIGHT, font=FONT_MONO)
    for im in range(1, 6):
        v = l * im
        cx = add(gx0, Emu((im - 1) * int(gcell)))
        cellc = band_color(v)
        cell = rect(s, cx, ry, gcell - Emu(int(Inches(0.03))), gcell - Emu(int(Inches(0.03))),
                    fill=cellc)
        txt = S_MED_INK if cellc == S_MED else WHITE
        setpar(cell.text_frame.paragraphs[0], str(v), Pt(13), color=txt, bold=True,
               align=PP_ALIGN.CENTER, font=FONT_MONO)

legend = [(S_LOW, "Low", "score 1–4"), (S_MED, "Medium", "score 5–9"),
          (S_HIGH, "High", "score 10–14"), (S_CRIT, "Critical", "score 15–25")]
lx = Inches(8.7)
ly = Inches(3.1)
for color, name, rng in legend:
    rect(s, lx, ly, Inches(0.22), Inches(0.22), fill=color)
    tb, tf = textbox(s, add(lx, Inches(0.35)), add(ly, Inches(-0.06)), Inches(3.5), Inches(0.35))
    add_runs(tf.paragraphs[0], [(name, INK, True), (f"  ·  {rng}", INK_SOFT, False)], Pt(12.5))
    ly = add(ly, Inches(0.5))

# =========================================================================
# S8. THE REGISTER — STAGE DISTRIBUTION BARS
# =========================================================================
s = base_slide(8)
kicker(s, "The register")
title(s, "Thirty risks, mapped across four stages", y=Inches(1.05), size=Pt(28))

bars = [("Use", 23), ("Development", 13), ("Deployment", 12), ("Monitoring", 6)]
maxv = max(v for _, v in bars)
bx0 = Inches(0.7)
blabel_w = Inches(2.0)
btrack_x = add(bx0, blabel_w)
btrack_w = Inches(8.0)
bval_x = add(btrack_x, add(btrack_w, Inches(0.15)))
by = Inches(2.4)
bh = Inches(0.5)
bgap = Inches(0.45)
for label, val in bars:
    tb, tf = textbox(s, bx0, by, blabel_w, bh, anchor=MSO_ANCHOR.MIDDLE)
    setpar(tf.paragraphs[0], label, Pt(13), color=INK_SOFT)
    rect(s, btrack_x, by, btrack_w, bh, fill=PAPER_RAISED, line=None)
    fillw = Emu(int(btrack_w * val / maxv))
    rect(s, btrack_x, by, fillw, bh, fill=ACCENT)
    tb2, tf2 = textbox(s, bval_x, by, Inches(0.6), bh, anchor=MSO_ANCHOR.MIDDLE)
    setpar(tf2.paragraphs[0], str(val), Pt(13), color=INK, font=FONT_MONO)
    by = add(by, add(bh, bgap))

footnote(s, "Counts sum to more than 30 because several risks are tagged to more than one "
            "lifecycle stage.", y=add(by, Inches(0.1)))

# =========================================================================
# S9. SPOTLIGHT
# =========================================================================
s = base_slide(9)
kicker(s, "Spotlight")
title(s, "A sample from the register", y=Inches(1.05), size=Pt(28))

spot = [
    ("AIR01", "Shadow AI Deployments",
     "AI tools or models deployed and used by internal teams without documentation, approval, "
     "or oversight from security or governance.", "DEPLOYMENT · USE"),
    ("AIR03", "Data Leakage",
     "Sensitive data exposed through third-party tools, API calls, or model interactions, "
     "leading to loss of confidentiality or regulatory breach.", "DEVELOPMENT · USE"),
    ("AIR15", "Adversarial Attacks",
     "An attacker deliberately alters input data to mislead the model — from prompt injection "
     "to imperceptible perturbations.", "USE"),
    ("AIR17", "Bias & Discrimination in Outputs",
     "Skewed or discriminatory results from biased training data or model design, including "
     "unequal performance across demographic groups.", "DEVELOPMENT · USE"),
    ("AIR27", "Unbounded Consumption",
     "Uncontrolled or excessive inference requests degrade availability or generate runaway "
     "compute costs — a.k.a. Denial of Wallet/Service.", "USE · MONITORING"),
    ("AIR29", "Multi-Agent & Agentic AI Risks",
     "Risks specific to autonomous agents acting with tools and memory via an agent harness — "
     "unconstrained tool access, unsafe actions, and privilege escalation.",
     "DEVELOPMENT · DEPLOYMENT · USE"),
]
cw = Inches(3.93)
ch = Inches(2.05)
cgap = Inches(0.015)
cx0 = Inches(0.7)
cy0 = Inches(2.1)
for i, (tag, name, body, stages) in enumerate(spot):
    col = i % 3
    row = i // 3
    x = add(cx0, Emu(col * int(add(cw, cgap))))
    y = add(cy0, Emu(row * int(add(ch, cgap))))
    rect(s, x, y, cw, ch, fill=PAPER, line=LINE, line_w=Pt(0.75))
    tb, tf = textbox(s, x + Inches(0.18), y + Inches(0.15), cw - Inches(0.36), ch - Inches(0.3))
    rect(s, x + Inches(0.18), y + Inches(0.15), Inches(0.7), Inches(0.28), fill=None,
         line=ACCENT, line_w=Pt(0.75))
    tbtag, tftag = textbox(s, x + Inches(0.18), y + Inches(0.15), Inches(0.9), Inches(0.28),
                           anchor=MSO_ANCHOR.MIDDLE)
    setpar(tftag.paragraphs[0], tag, Pt(9), color=ACCENT, font=FONT_MONO)
    tb2, tf2 = textbox(s, x + Inches(0.18), y + Inches(0.5), cw - Inches(0.36), Inches(0.4))
    setpar(tf2.paragraphs[0], name, Pt(12.5), color=INK, font=FONT_HEAD, line=1.1)
    tb3, tf3 = textbox(s, x + Inches(0.18), y + Inches(0.95), cw - Inches(0.36), Inches(0.85))
    setpar(tf3.paragraphs[0], body, Pt(9.5), color=INK_SOFT, line=1.4)
    tb4, tf4 = textbox(s, x + Inches(0.18), y + ch - Inches(0.32), cw - Inches(0.36), Inches(0.28))
    setpar(tf4.paragraphs[0], stages, Pt(8), color=INK_FAINT, font=FONT_MONO)

# =========================================================================
# S10. FULL REGISTER
# =========================================================================
register = [
    ("AIR01", "Shadow AI Deployments", "Deployment, Use"),
    ("AIR02", "Supply Chain Vulnerabilities", "Development, Deployment"),
    ("AIR03", "Data Leakage", "Development, Use"),
    ("AIR04", "Cross-Border Data Transfer", "Deployment, Use"),
    ("AIR05", "No Validation / Improper Output", "Use"),
    ("AIR06", "Insecure Deployment Pipelines", "Deployment"),
    ("AIR07", "Critical 3rd-party Dependency", "Development, Deployment"),
    ("AIR08", "Excessive Agency / Over-Reliance on AI", "Use"),
    ("AIR09", "Model Poisoning", "Development"),
    ("AIR10", "Model Inversion / Membership Inference / Theft", "Deployment, Use"),
    ("AIR11", "Poor Monitoring", "Monitoring"),
    ("AIR12", "Feedback Loop Contamination", "Use, Monitoring"),
    ("AIR13", "Alert Fatigue", "Monitoring"),
    ("AIR14", "Insider Threats", "Development, Deployment, Use"),
    ("AIR15", "Adversarial Attacks", "Use"),
    ("AIR16", "Model Drift", "Use, Monitoring"),
    ("AIR17", "Bias & Discrimination in Outputs", "Development, Use"),
    ("AIR18", "Critical Unintended Consequences", "Development, Deployment, Use"),
    ("AIR19", "Censorship/Guardrail interference", "Use"),
    ("AIR20", "Deployment misfit", "Deployment, Use"),
    ("AIR21", "Lack of Explainability / Provenance opacity", "Development, Use"),
    ("AIR22", "No Accountability", "Development, Deployment, Monitoring"),
    ("AIR23", "Unclear Model Ownership/IP", "Development, Use"),
    ("AIR24", "Output Integrity Tampering", "Use"),
    ("AIR25", "System Prompt / Instruction Leakage", "Use"),
    ("AIR26", "RAG / Vector Store Vulnerabilities", "Development, Use"),
    ("AIR27", "Unbounded Consumption", "Use, Monitoring"),
    ("AIR28", "Malicious Use / Out-of-Scope Exploitation", "Deployment, Use"),
    ("AIR29", "Multi-Agent & Agentic AI Risks", "Development, Deployment, Use"),
    ("AIR30", "Environmental / Sustainability Impact", "Development, Use"),
]

s = base_slide(10)
kicker(s, "Full register")
title(s, "All thirty, AIR01–AIR30", y=Inches(1.05), size=Pt(26))

row_h = Inches(0.315)
gx = Inches(0.7)
gy = Inches(1.9)
half = 15
for half_i, items in enumerate([register[:half], register[half:]]):
    x0 = add(gx, Emu(half_i * int(Inches(6.35))))
    hh = Inches(0.32)
    tb, tf = textbox(s, x0, gy, Inches(6.15), hh, anchor=MSO_ANCHOR.MIDDLE)
    add_runs(tf.paragraphs[0], [("ID", INK_SOFT, False), ("      ", INK_SOFT, False),
                                 ("RISK", INK_SOFT, False)], Pt(9), font=FONT_MONO)
    ry = add(gy, hh)
    for i, (aid, name, stg) in enumerate(items):
        fillc = PAPER_RAISED if i % 2 == 0 else PAPER
        rect(s, x0, ry, Inches(6.15), row_h, fill=fillc, line=LINE, line_w=Pt(0.4))
        tb, tf = textbox(s, x0 + Inches(0.08), ry, Inches(0.85), row_h, anchor=MSO_ANCHOR.MIDDLE)
        setpar(tf.paragraphs[0], aid, Pt(9), color=ACCENT, font=FONT_MONO)
        tb2, tf2 = textbox(s, x0 + Inches(0.9), ry, Inches(3.85), row_h, anchor=MSO_ANCHOR.MIDDLE)
        setpar(tf2.paragraphs[0], name, Pt(9), color=INK, line=1.0)
        tb3, tf3 = textbox(s, x0 + Inches(4.7), ry, Inches(1.4), row_h, anchor=MSO_ANCHOR.MIDDLE)
        setpar(tf3.paragraphs[0], stg, Pt(7.5), color=INK_FAINT, line=1.0)
        ry = add(ry, row_h)

# =========================================================================
# S11. STANDARDS MAPPING
# =========================================================================
s = base_slide(11)
kicker(s, "Standards")
title(s, "Mapped to existing governance", y=Inches(1.05), size=Pt(28))
tb, tf = textbox(s, Inches(0.7), Inches(1.85), Inches(11.4), Inches(0.6))
setpar(tf.paragraphs[0],
       "Every AIR## risk is cross-referenced so findings plug straight into the programs you "
       "already run.", Pt(14.5), color=INK_SOFT, font=FONT_SUB, line=1.4)

chips = [
    ("NIST AI RMF", "Govern · Map · Measure · Manage"),
    ("OWASP LLM Top 10", "2025 edition"),
    ("EU AI Act", "Articles 5–95"),
    ("ISO/IEC 42001", "Annex A controls"),
    ("MITRE ATLAS", "AML.T#### techniques"),
]
cx = Inches(0.7)
chw = Inches(2.34)
chg = Inches(0.06)
chy = Inches(2.55)
for name, sub in chips:
    rect(s, cx, chy, chw, Inches(0.8), fill=PAPER_RAISED, line=LINE, line_w=Pt(0.75))
    tb, tf = textbox(s, cx + Inches(0.12), chy + Inches(0.1), chw - Inches(0.24), Inches(0.6))
    setpar(tf.paragraphs[0], name, Pt(11.5), color=INK, font=FONT_HEAD, space_after=Pt(2))
    p2 = tf.add_paragraph()
    setpar(p2, sub, Pt(8.5), color=INK_FAINT, font=FONT_MONO, line=1.1)
    cx = add(add(cx, chw), chg)

map_rows = [
    ("AIR", "Risk", "OWASP LLM", "EU AI Act", "MITRE ATLAS"),
    ("AIR03", "Data Leakage", "LLM02 Sensitive Info", "Art 10", "AML.T0057"),
    ("AIR15", "Adversarial Attacks", "LLM01 Prompt Injection", "Art 15", "AML.T0051"),
    ("AIR26", "RAG / Vector Store", "LLM08 Vector Weakness", "Art 10 & 15", "AML.T0070"),
    ("AIR27", "Unbounded Consumption", "LLM10 Unbounded", "Art 15", "AML.T0029"),
]
end_y = table_grid(s, Inches(0.7), Inches(3.65), [Inches(1.1), Inches(2.6), Inches(2.9),
                    Inches(1.7), Inches(1.6)], map_rows, row_h=Inches(0.48), font_size=Pt(10.5))
footnote(s, "Full mapping in standards-map.csv.", y=add(end_y, Inches(0.2)))

# =========================================================================
# S12. FROM REGISTER TO RISK LOG
# =========================================================================
s = base_slide(12)
kicker(s, "In practice")
title(s, "From register to risk log", y=Inches(1.05), size=Pt(28))
tb, tf = textbox(s, Inches(0.7), Inches(1.85), Inches(11.4), Inches(0.7))
setpar(tf.paragraphs[0],
       "Record two ratings, an owner, and a control per applicable risk in your project’s "
       "log — traceable back to the register by AIR ID.", Pt(14.5), color=INK_SOFT,
       font=FONT_SUB, line=1.4)

log_rows = [
    ("Initiative", "AIR", "Risk", "L", "I", "Score", "Level", "Owner", "Status"),
    ("Support chatbot", "AIR05", "No Validation / Improper Output", "3", "4", "12", "High",
     "ML Lead", "Mitigating"),
    ("Support chatbot", "AIR27", "Unbounded Consumption", "2", "3", "6", "Medium", "Platform",
     "Accepted"),
]
log_ws = [Inches(1.7), Inches(0.85), Inches(2.9), Inches(0.4), Inches(0.4), Inches(0.65),
          Inches(0.85), Inches(1.0), Inches(1.1)]
end_y = table_grid(s, Inches(0.7), Inches(2.75), log_ws, log_rows, row_h=Inches(0.55),
                    font_size=Pt(10))
footnote(s, "Template fields: likelihood, impact, score, level, controls, owner, sign-off, "
            "KRI / monitoring signal, next review, status — see risk-log.template.csv.",
         y=add(end_y, Inches(0.2)))

# =========================================================================
# S13. WHAT'S INSIDE (REPO ASSETS)
# =========================================================================
s = base_slide(13)
kicker(s, "What’s inside")
title(s, "Four files, ready to use", y=Inches(1.05), size=Pt(28))

assets = [
    ("register.csv", "The 30-risk AI Risk Register with IDs, lifecycle tags, and descriptions."),
    ("standards-map.csv", "Cross-references to NIST, OWASP, EU AI Act, ISO 42001 & MITRE ATLAS."),
    ("risk-log.template.csv", "Ready-to-fill log for scoring, ownership, sign-off, and review cadence."),
    ("README.md", "Lifecycle model, usage workflow, and scoring guidance."),
]
acw = Inches(5.93)
ach = Inches(1.55)
acgap = Inches(0.015)
ax0 = Inches(0.7)
ay0 = Inches(2.3)
for i, (fname, body) in enumerate(assets):
    col = i % 2
    row = i // 2
    x = add(ax0, Emu(col * int(add(acw, acgap))))
    y = add(ay0, Emu(row * int(add(ach, acgap))))
    rect(s, x, y, acw, ach, fill=PAPER, line=LINE, line_w=Pt(0.75))
    tb, tf = textbox(s, x + Inches(0.2), y + Inches(0.18), acw - Inches(0.4), ach - Inches(0.3))
    setpar(tf.paragraphs[0], fname, Pt(13.5), color=ACCENT, font=FONT_MONO, space_after=Pt(6))
    p2 = tf.add_paragraph()
    setpar(p2, body, Pt(11), color=INK_SOFT, line=1.4)

# =========================================================================
# S14. CLOSING
# =========================================================================
s = slide()
bg(s, DARK_BG)
kicker_tb, kicker_tf = textbox(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.35))
setpar(kicker_tf.paragraphs[0], "GET INVOLVED", Pt(11), color=RGBColor(0x8A, 0xC7, 0x5A), bold=True,
       font=FONT_HEAD)

tb, tf = textbox(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(1.4))
setpar(tf.paragraphs[0], "Treat it as a living list.", Pt(42), color=WHITE, bold=False,
       font=FONT_HEAD)

tb, tf = textbox(s, Inches(0.7), Inches(3.35), Inches(9.2), Inches(1.1))
setpar(tf.paragraphs[0],
       "If you identify a risk not represented here, propose it by opening an issue in the "
       "repository — your contribution may help other teams facing the same gap.",
       Pt(15.5), color=RGBColor(0xDD, 0xDD, 0xDD), font=FONT_SUB, line=1.5)

cta_x = Inches(0.7)
cta_y = Inches(4.85)
cta_w = Inches(4.85)
cta_h = Inches(1.05)
cta = [("REPOSITORY", "github.com/OpenEthicsAI/RSAIRF"), ("LICENSE", "CC BY 4.0")]
for k, v in cta:
    rect(s, cta_x, cta_y, cta_w, cta_h, fill=None, line=RGBColor(0x34, 0x38, 0x3C), line_w=Pt(0.75))
    tb, tf = textbox(s, cta_x + Inches(0.2), cta_y + Inches(0.15), cta_w - Inches(0.4),
                     cta_h - Inches(0.25))
    setpar(tf.paragraphs[0], k, Pt(9.5), color=RGBColor(0x8A, 0x8A, 0x8A), bold=True,
           font=FONT_HEAD)
    p2 = tf.add_paragraph()
    setpar(p2, v, Pt(13.5), color=WHITE, space_before=Pt(6))
    cta_x = add(cta_x, add(cta_w, Emu(int(Inches(0.02)))))

# ---- save ----------------------------------------------------------------
prs.save("RSAIRF.pptx")
print(f"Saved RSAIRF.pptx with {len(prs.slides._sldIdLst)} slides")
